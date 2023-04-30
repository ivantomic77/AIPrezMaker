import os
import openai
from dotenv import load_dotenv
from pptx import Presentation
import json

load_dotenv()

def main():
  openai.organization = os.getenv('ORG')
  openai.api_key = os.getenv('API_KEY')
  index = 0;

  language = input("What language should the presentation use?\n")
  theme = input("What is the topic for presentation?\n")
  minSlides = input("How much slides should it contain?\n")

  prompt = 'Make me a complex presentation in language: %s. The required number of slides is: %s. The theme is: %s.  Presentation should not have a lot of text on one slide. Give me a presentation in this format strictly: { "Slide number": {"title": "Title", "text": "Slide text", "prompt": "Prompt I can feed to image generation" }, ...}' % (language, minSlides,theme)
  result = openai.Completion.create(
    model="text-davinci-003",
    prompt=prompt,
    max_tokens=1000,
    temperature=1
  )
  pptInJSON = json.loads(result.choices[0].text)

  presentation = Presentation()
  firstLayout = presentation.slide_layouts[0]  
  mySlide = presentation.slides.add_slide(firstLayout)
  myTitle = mySlide.shapes.title
  mySubtitle = mySlide.shapes.placeholders[1]
  myTitle.text = theme

  for key in pptInJSON:
    value = pptInJSON[key]
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = value["title"]
    body = slide.shapes.placeholders[1].text_frame
    body.text = value["text"]
  
  presentation.save('myPPT.pptx')
    

if __name__ == '__main__':
    main()
